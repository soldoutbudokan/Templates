# %%
################################################################################
# Title: Document Review - PPTX
# Author: Tirth Bhatt
#
# Project: [PROJECT]
################################################################################
# Interactive PPTX Keyword Search
# Run in terminal: python pptx_search_interactive.py

import os
from pathlib import Path
from pptx import Presentation
import re
import csv

# %%
# CONFIGURATION - Edit this path

INPUT_FOLDER = r"[PATH]"
CONTEXT_CHARS = 150

# %%
def load_pptx_files(folder_path):
    """
    Extract text from all PPTX files upfront. Returns dict of {filename: {slide_num: text}}
    """
    folder = Path(folder_path)
    pptx_files = list(folder.rglob("*.pptx"))
    
    print(f"Loading {len(pptx_files)} PPTX files...")
    
    pptx_texts = {}
    
    for i, pptx_path in enumerate(pptx_files, start=1):
        filename = pptx_path.name
        print(f"  [{i}/{len(pptx_files)}] {filename}", end="\r")
        
        try:
            prs = Presentation(pptx_path)
            slides = {}
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                if slide_text:
                    combined = " ".join(slide_text)
                    slides[slide_num] = re.sub(r'\s+', ' ', combined)
            pptx_texts[filename] = slides
        except Exception as e:
            print(f"\n  Error loading {filename}: {e}")
    
    print(f"\nLoaded {len(pptx_texts)} PPTX files successfully.")
    return pptx_texts


# %%
def search_keyword(pptx_texts, keyword, context_chars=150):
    """
    Search all loaded PPTX files for a keyword.
    """
    results = []
    pattern = re.compile(re.escape(keyword), re.IGNORECASE)
    
    for filename, slides in pptx_texts.items():
        for slide_num, text in slides.items():
            for match in pattern.finditer(text):
                start = max(0, match.start() - context_chars)
                end = min(len(text), match.end() + context_chars)
                
                context = text[start:end]
                if start > 0:
                    context = "..." + context
                if end < len(text):
                    context = context + "..."
                
                results.append({
                    'file': filename,
                    'slide': slide_num,
                    'context': context.strip()
                })
    
    return results


# %%
def display_results(results, keyword):
    """
    Print results in a readable format.
    """
    if not results:
        print(f"\nNo matches for '{keyword}'")
        return
    
    print(f"\n{'='*60}")
    print(f"Found {len(results)} matches for '{keyword}'")
    print('='*60)
    
    for r in results:
        print(f"\n[{r['file']}] Slide {r['slide']}")
        print(f"  {r['context']}")
    
    # Summary
    files = set(r['file'] for r in results)
    print(f"\n--- Summary: {len(results)} matches across {len(files)} files ---")


# %%
def export_results(results, keyword):
    """
    Save results to CSV.
    """
    filename = f"search_{keyword.replace(' ', '_')}.csv"
    with open(filename, 'w', newline='', encoding='utf-8') as f:
        writer = csv.DictWriter(f, fieldnames=['file', 'slide', 'context'])
        writer.writeheader()
        writer.writerows(results)
    print(f"Saved to: {filename}")


# %%
def main():
    print("="*60)
    print("PPTX KEYWORD SEARCH")
    print("="*60)
    
    # Load PPTX files once
    pptx_texts = load_pptx_files(INPUT_FOLDER)
    
    if not pptx_texts:
        print("No PPTX files found. Check your INPUT_FOLDER path.")
        return
    
    last_results = []
    last_keyword = ""
    
    print("\nCommands:")
    print("  Type a keyword to search")
    print("  'export' - save last results to CSV")
    print("  'quit' or 'exit' - close program")
    print()
    
    while True:
        try:
            query = input("Search> ").strip()
        except (KeyboardInterrupt, EOFError):
            print("\nGoodbye.")
            break
        
        if not query:
            continue
        
        if query.lower() in ('quit', 'exit', 'q'):
            print("Goodbye.")
            break
        
        if query.lower() == 'export':
            if last_results:
                export_results(last_results, last_keyword)
            else:
                print("No results to export. Run a search first.")
            continue
        
        # Run search
        last_keyword = query
        last_results = search_keyword(pptx_texts, query, CONTEXT_CHARS)
        display_results(last_results, query)


# %%
if __name__ == "__main__":
    main()
